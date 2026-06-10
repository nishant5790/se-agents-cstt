"""Agent 4 — PPT Generation. Renders a DeckPlan into a .pptx using python-pptx.

`python-pptx` is the agent's "tool". Slide layouts are small functions so new
slide styles are easy to add. Uses templates/brand.pptx as the master if present.
"""
from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from agent_team.core.base import Agent, Blackboard
from .analysis_agent import DeckPlan, SlidePlan


def _safe(name: str) -> str:
    return re.sub(r"[^\w\- ]+", "", name).strip()[:80] or "deck"


class PptAgent(Agent):
    name = "ppt"

    def __init__(self, out_dir: Path, template: Path | None = None):
        self.out_dir = out_dir
        self.template = template if template and template.exists() else None

    def run(self, bb: Blackboard) -> Blackboard:
        plan: DeckPlan = bb.get("deck_plan")
        if plan is None:
            raise SystemExit("Run AnalysisAgent first.")
        path = self.build(plan)
        bb.set("pptx_path", str(path))
        self.log(f"wrote {path.name} ({len(plan.slides)} slides)")
        return bb

    # --- tool: build a pptx ---
    def build(self, plan: DeckPlan) -> Path:
        prs = Presentation(str(self.template)) if self.template else Presentation()
        for i, slide in enumerate(plan.slides):
            if i == 0:
                self._title_slide(prs, plan.deck_title, slide)
            else:
                self._content_slide(prs, slide)
        self.out_dir.mkdir(parents=True, exist_ok=True)
        out = self.out_dir / f"DECK--{_safe(plan.deck_title)}.pptx"
        prs.save(str(out))
        return out

    def _title_slide(self, prs: Presentation, deck_title: str, slide: SlidePlan) -> None:
        layout = prs.slide_layouts[0]
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = deck_title
        if len(s.placeholders) > 1:
            s.placeholders[1].text = " | ".join(slide.bullets) if slide.bullets else slide.title

    def _content_slide(self, prs: Presentation, slide: SlidePlan) -> None:
        layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        s = prs.slides.add_slide(layout)
        if s.shapes.title is not None:
            s.shapes.title.text = slide.title
        has_image = bool(slide.image and Path(slide.image).exists())
        body = self._body_placeholder(s)
        if body is not None:
            if has_image:
                # narrow the text column to the left half to make room for the image
                try:
                    body.left, body.top = Inches(0.6), Inches(1.6)
                    body.width, body.height = Inches(4.7), Inches(5.0)
                except Exception:
                    pass
            tf = body.text_frame
            tf.word_wrap = True
            tf.clear()
            for j, bullet in enumerate(slide.bullets or [slide.title]):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.text = bullet
                p.level = 0
                for run in p.runs:
                    run.font.size = Pt(18)
        if has_image:
            self._place_image(s, slide.image)
        if slide.notes:
            s.notes_slide.notes_text_frame.text = slide.notes

    @staticmethod
    def _place_image(slide, image_path: str) -> None:
        """Place an image on the right half of the slide, scaled to fit."""
        try:
            slide.shapes.add_picture(
                image_path, Inches(5.5), Inches(1.6), height=Inches(4.5))
        except Exception:
            pass

    @staticmethod
    def _body_placeholder(slide):
        for ph in slide.placeholders:
            if ph.placeholder_format.idx != 0:  # not the title
                return ph
        # fall back to a fresh textbox
        return slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(8.5), Inches(5))
